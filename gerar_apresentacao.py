from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io

# ── Paleta de cores ──────────────────────────────────────────────────────────
AMARELO      = RGBColor(0xFF, 0xC1, 0x07)   # amarelo sol
AMARELO_ESC  = RGBColor(0xE6, 0xA8, 0x00)
AZUL_MAR     = RGBColor(0x00, 0x6A, 0xA7)   # azul mar
AZUL_CLA     = RGBColor(0xE8, 0xF4, 0xFB)   # fundo claro
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_ESC    = RGBColor(0x2C, 0x2C, 0x2C)
CINZA_MED    = RGBColor(0x55, 0x55, 0x55)
VERDE        = RGBColor(0x28, 0xA7, 0x45)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout completamente em branco


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=CINZA_ESC,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_text_multiline(slide, lines, l, t, w, h,
                       font_size=16, bold=False, color=CINZA_ESC,
                       align=PP_ALIGN.LEFT, line_spacing=1.15):
    """lines = list of str; first item can be a tuple (text, bold, color)"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple):
            txt, b, c = line
        else:
            txt, b, c = line, bold, color
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = b
        run.font.color.rgb = c


def bullet_block(slide, items, l, t, w, h,
                 font_size=16, color=CINZA_ESC, bullet="●  "):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def slide_header(slide, title, subtitle=None):
    """Faixa azul no topo com título e subtítulo opcionais."""
    add_rect(slide, 0, 0, 13.33, 1.35, AZUL_MAR)
    add_rect(slide, 0, 1.35, 13.33, 0.08, AMARELO)
    add_text(slide, title, 0.45, 0.15, 12.4, 0.75,
             font_size=28, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.45, 0.82, 12.4, 0.45,
                 font_size=15, bold=False, color=AMARELO, align=PP_ALIGN.LEFT)


def card(slide, l, t, w, h, bg=AZUL_CLA, border=AZUL_MAR):
    add_rect(slide, l, t, w, h, bg, border)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 – Capa
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, AZUL_MAR)
add_rect(sl, 0, 5.5, 13.33, 2.0, AMARELO)

# Faixa diagonal decorativa (retângulo inclinado simulado com dois retângulos)
add_rect(sl, 0, 3.8, 13.33, 0.12, AMARELO_ESC)

# Título principal
add_text(sl, "☀ Sunny Sales", 1.0, 1.2, 11.0, 1.2,
         font_size=54, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)

add_text(sl, "Comércio de Praia Inteligente", 1.0, 2.45, 11.0, 0.7,
         font_size=28, bold=False, color=BRANCO, align=PP_ALIGN.CENTER)

add_text(sl, "Apresentação para Câmara Municipal", 1.0, 3.2, 11.0, 0.5,
         font_size=17, bold=False, color=BRANCO, align=PP_ALIGN.CENTER)

add_text(sl, "Solução digital para organizar e modernizar os vendedores ambulantes de praia",
         1.0, 5.7, 11.0, 0.6,
         font_size=16, bold=False, color=AZUL_MAR, align=PP_ALIGN.CENTER)

add_text(sl, "suporte@sunnysales.com", 5.5, 6.6, 3.0, 0.4,
         font_size=13, bold=False, color=AZUL_MAR, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 – Índice
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Agenda", "O que vamos apresentar hoje")

topics = [
    ("01", "O Problema", "Desafios actuais do comércio de praia"),
    ("02", "A Solução", "Sunny Sales — como funciona"),
    ("03", "Para os Banhistas", "Localizar vendedores em tempo real"),
    ("04", "Para os Vendedores", "Painel, rotas e subscrição"),
    ("05", "Para o Município", "Controlo, dados e implementação"),
    ("06", "Implementação", "Como activar em dias"),
    ("07", "Sustentabilidade", "Impacto ambiental positivo"),
    ("08", "Próximos Passos", "Como avançar"),
]

cols = [(0.4, 1.65), (6.9, 1.65), (0.4, 3.4), (6.9, 3.4),
        (0.4, 5.15), (6.9, 5.15), (0.4, 6.7), (6.9, 6.7)]

for idx, (num, tit, sub) in enumerate(topics):
    lx, ty = cols[idx]
    card(sl, lx, ty, 6.2, 1.5, AZUL_CLA)
    add_text(sl, num, lx + 0.12, ty + 0.12, 0.6, 0.55,
             font_size=22, bold=True, color=AMARELO_ESC)
    add_text(sl, tit, lx + 0.75, ty + 0.12, 5.2, 0.45,
             font_size=17, bold=True, color=AZUL_MAR)
    add_text(sl, sub, lx + 0.75, ty + 0.58, 5.2, 0.65,
             font_size=13, color=CINZA_MED)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 – O Problema
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "O Problema", "O caos invisível no comércio de praia")

problems = [
    ("🔍  Banhistas perdem tempo", "Percorrem grandes distâncias à procura de gelados, pastéis ou acessórios — sem saber onde o vendedor está."),
    ("🚶  Vendedores andam em excesso", "Percorrem a praia aleatoriamente, sem informação sobre a localização dos potenciais clientes."),
    ("📉  Comércio desorganizado", "Sobreposição de rotas, zonas sem cobertura, ausência de dados para a autarquia gerir concessões."),
    ("❌  Sem visibilidade digital", "Os vendedores ambulantes não têm presença digital — invisíveis para turistas e visitantes."),
]

for i, (title, desc) in enumerate(problems):
    row = i // 2
    col = i % 2
    lx = 0.4 + col * 6.45
    ty = 1.65 + row * 2.6
    add_rect(sl, lx, ty, 6.2, 2.35,
             RGBColor(0xFF, 0xF3, 0xCD) if row == 0 and col == 0 else
             RGBColor(0xFF, 0xDD, 0xDD) if row == 0 else
             RGBColor(0xDD, 0xF0, 0xFF) if col == 0 else
             RGBColor(0xE8, 0xF5, 0xE9),
             None)
    add_text(sl, title, lx + 0.2, ty + 0.18, 5.8, 0.55,
             font_size=16, bold=True, color=CINZA_ESC)
    add_text(sl, desc, lx + 0.2, ty + 0.75, 5.8, 1.4,
             font_size=14, color=CINZA_MED)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 – A Solução
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "A Solução", "Sunny Sales — Plataforma de Comércio de Praia Inteligente")

add_text(sl,
    "O Sunny Sales liga banhistas e vendedores em tempo real através de GPS e mapas interactivos, "
    "eliminando a incerteza de ambos os lados.",
    0.5, 1.55, 12.3, 0.8, font_size=16, color=CINZA_ESC)

# Três pilares
for i, (icon, lbl, desc) in enumerate([
    ("📱", "Sem instalar app", "Acesso instantâneo por QR Code — o banhista abre o mapa no browser"),
    ("📍", "Localização em tempo real", "O vendedor partilha a sua posição GPS; o mapa actualiza automaticamente via WebSocket"),
    ("📊", "Dados para o município", "Estatísticas de actividade, percursos, cobertura e utilização em tempo real"),
]):
    lx = 0.5 + i * 4.25
    add_rect(sl, lx, 2.55, 3.95, 4.5, AZUL_MAR)
    add_text(sl, icon, lx, 2.75, 3.95, 0.8, font_size=38, align=PP_ALIGN.CENTER, color=BRANCO)
    add_text(sl, lbl, lx + 0.15, 3.55, 3.65, 0.6,
             font_size=16, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
    add_text(sl, desc, lx + 0.2, 4.2, 3.55, 2.5,
             font_size=13, color=BRANCO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 – Para os Banhistas
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Para os Banhistas", "Encontrar o vendedor em segundos, sem sair da toalha")

add_rect(sl, 0, 1.43, 13.33, 6.07, AZUL_CLA)

features = [
    ("🗺️  Mapa interactivo",     "Pinos coloridos e actualizados ao vivo mostram todos os vendedores activos na praia."),
    ("📡  GPS próprio",           "A posição do banhista é mostrada no mapa com seta de direcção (bússola do telemóvel)."),
    ("🔎  Filtros inteligentes",  "Filtrar por tipo de produto (gelados, pastéis, acessórios) e por distância (até 5 km)."),
    ("📸  Perfil do vendedor",    "Foto, produto e histórias efémeras (fotos/vídeos temporários) no perfil do vendedor."),
    ("📲  Sem instalação",        "Funciona directamente no browser via QR Code — sem app store, sem registo obrigatório."),
    ("🔄  Actualizações live",    "WebSocket garante actualizações automáticas — sem necessidade de refrescar a página."),
]

for i, (tit, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    lx = 0.4 + col * 6.45
    ty = 1.65 + row * 1.85
    add_rect(sl, lx, ty, 6.2, 1.7, BRANCO)
    add_text(sl, tit, lx + 0.2, ty + 0.12, 5.8, 0.5,
             font_size=15, bold=True, color=AZUL_MAR)
    add_text(sl, desc, lx + 0.2, ty + 0.62, 5.8, 0.9,
             font_size=13, color=CINZA_MED)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 – Para os Vendedores
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Para os Vendedores", "Ferramenta profissional de gestão e visibilidade")

cols_data = [
    ("Painel Pessoal", [
        "Dashboard com saudação personalizada",
        "Toggle para activar/desactivar partilha de localização",
        "Estado da subscrição e dias restantes",
        "Histórico de sessões em múltiplos dispositivos",
    ]),
    ("Rotas & Estatísticas", [
        "Registo automático de cada sessão de trabalho",
        "Distância percorrida e duração por sessão",
        "Mapa visual de cada rota efectuada",
        "Gráfico de barras de km diários (últimos 7 dias)",
    ]),
    ("Perfil & Conta", [
        "Foto de perfil com ferramenta de recorte",
        "Cor personalizada do pino no mapa",
        "Histórias efémeras (fotos/vídeos temporários)",
        "Gestão de password e dados pessoais",
    ]),
]

for i, (title, items) in enumerate(cols_data):
    lx = 0.4 + i * 4.25
    add_rect(sl, lx, 1.6, 3.95, 5.6, AZUL_CLA)
    add_rect(sl, lx, 1.6, 3.95, 0.65, AZUL_MAR)
    add_text(sl, title, lx + 0.15, 1.68, 3.65, 0.5,
             font_size=15, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(sl, "✔  " + item, lx + 0.2, 2.45 + j * 1.1, 3.55, 0.95,
                 font_size=13, color=CINZA_ESC)

add_text(sl, "💳  Pagamento de subscrição integrado com Stripe — simples, seguro e automático.",
         0.5, 7.1, 12.3, 0.38, font_size=13, bold=True, color=AZUL_MAR)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 – Para o Município
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Para o Município", "Controlo, transparência e modernização")

add_rect(sl, 0, 1.43, 13.33, 6.07, RGBColor(0xF0, 0xF7, 0xFD))

benefits = [
    ("📋  Gestão de concessões",
     "Registo digital de todos os vendedores licenciados. Subscrição activa = vendedor autorizado. Visibilidade total sobre quem está activo."),
    ("📊  Estatísticas em tempo real",
     "Dados de utilização: quantos vendedores activos, zonas de maior cobertura, frequência de uso por dia e hora."),
    ("🏖️  Melhoria da experiência do visitante",
     "Praia mais organizada e agradável. Turistas e famílias encontram facilmente o que precisam, aumentando a satisfação."),
    ("🌱  Compromisso de sustentabilidade",
     "Redução de percursos desnecessários → menos emissões. Dados de distâncias percorridas pelas vendedores disponíveis."),
    ("⚡  Implementação ultra-rápida",
     "QR Codes impressos → sistema operacional em dias. Sem infraestrutura adicional. Sem custos de hardware."),
    ("🤝  Suporte técnico dedicado",
     "Equipa disponível para integração, formação e personalização. Página web personalizada com o nome da praia/município."),
]

for i, (tit, desc) in enumerate(benefits):
    row = i // 2
    col = i % 2
    lx = 0.4 + col * 6.45
    ty = 1.65 + row * 1.9
    add_rect(sl, lx, ty, 6.2, 1.75, BRANCO)
    add_text(sl, tit, lx + 0.2, ty + 0.1, 5.8, 0.5,
             font_size=15, bold=True, color=AZUL_MAR)
    add_text(sl, desc, lx + 0.2, ty + 0.6, 5.8, 1.0,
             font_size=12.5, color=CINZA_MED)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 – Como Funciona (fluxo)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Como Funciona", "Três passos simples — do QR Code ao mapa")

steps = [
    ("1", "🏛️\nMunicípio\ncoloca QR Codes\nna praia",     AZUL_MAR),
    ("→", "",                                              BRANCO),
    ("2", "🧑‍🤝‍🧑\nBanhista\ndifusão o QR Code\nno telemóvel",  RGBColor(0x00, 0x8C, 0x60)),
    ("→", "",                                              BRANCO),
    ("3", "🗺️\nMapa abre\nno browser —\nsem app",          RGBColor(0x6F, 0x42, 0xC1)),
    ("→", "",                                              BRANCO),
    ("4", "📍\nVendedor\npartilha GPS\nem tempo real",      RGBColor(0xDC, 0x35, 0x45)),
]

positions = [0.3, 2.55, 3.0, 5.25, 5.7, 7.95, 8.4]
for idx, (num, lbl, color) in enumerate(steps):
    lx = positions[idx]
    if num == "→":
        add_text(sl, "→", lx, 3.3, 0.6, 0.8,
                 font_size=32, bold=True, color=RGBColor(0xAA, 0xAA, 0xAA),
                 align=PP_ALIGN.CENTER)
    else:
        add_rect(sl, lx, 2.0, 2.35, 4.0, color)
        add_text(sl, lbl, lx + 0.1, 2.15, 2.15, 3.7,
                 font_size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

add_rect(sl, 0.3, 6.2, 12.73, 1.0, AZUL_CLA)
add_text(sl,
    "O ciclo fecha-se: o banhista vê o vendedor no mapa → desloca-se → compra. "
    "Município recebe dados de toda a actividade.",
    0.6, 6.3, 12.1, 0.75, font_size=14, color=AZUL_MAR, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 – Implementação
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Implementação", "Do contrato à praia operacional em poucos dias")

timeline = [
    ("Dia 1–2",  "Contrato & Configuração",
     "Definição das praias, número de licenças, personalização do sistema com o nome do município."),
    ("Dia 3–4",  "Onboarding de Vendedores",
     "Registo dos vendedores licenciados, instalação da app mobile (opcional), formação básica."),
    ("Dia 5",    "QR Codes & Lançamento",
     "Impressão e colocação dos QR Codes nos postos de praia, bandeiras e infraestruturas existentes."),
    ("Dia 6+",   "Operação & Suporte",
     "Sistema em funcionamento. Equipa técnica disponível. Relatórios mensais ao município."),
]

for i, (period, title, desc) in enumerate(timeline):
    ty = 1.7 + i * 1.4
    add_rect(sl, 0.4, ty, 1.8, 1.2, AZUL_MAR)
    add_text(sl, period, 0.4, ty + 0.05, 1.8, 1.1,
             font_size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.3, ty, 10.6, 1.2, AZUL_CLA)
    add_text(sl, title, 2.5, ty + 0.05, 10.2, 0.45,
             font_size=16, bold=True, color=AZUL_MAR)
    add_text(sl, desc, 2.5, ty + 0.5, 10.2, 0.6,
             font_size=13, color=CINZA_MED)

add_rect(sl, 0.4, 7.1, 12.5, 0.35, AMARELO)
add_text(sl, "🚀  Requisito mínimo: ligação à internet. Sem hardware dedicado. Sem obras.",
         0.6, 7.12, 12.1, 0.3, font_size=14, bold=True, color=AZUL_MAR,
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 – Sustentabilidade
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Sustentabilidade", "Um passo para praias mais inteligentes e sustentáveis")

add_rect(sl, 0, 1.43, 13.33, 6.07, RGBColor(0xF0, 0xFB, 0xF4))

points = [
    ("🌍  Menos quilómetros, menos CO₂",
     "Com vendedores a dirigirem-se para onde há procura, eliminam-se percursos desnecessários. "
     "A plataforma regista e apresenta os km poupados."),
    ("♻️  Sem papel, sem plástico",
     "Toda a comunicação é digital. QR Codes substituem panfletos e mapas em papel. "
     "Menor pegada ecológica no dia-a-dia da praia."),
    ("📈  Dados de impacto ambiental",
     "O município pode aceder a relatórios de distâncias percorridas e quantificar "
     "a redução de emissões de carbono da sua praia."),
    ("🏖️  Praia mais organizada",
     "Menos conflitos de sobreposição de rotas. Melhor distribuição dos vendedores. "
     "Experiência mais agradável para visitantes e residentes."),
]

for i, (tit, desc) in enumerate(points):
    row = i // 2
    col = i % 2
    lx = 0.4 + col * 6.45
    ty = 1.7 + row * 2.55
    add_rect(sl, lx, ty, 6.2, 2.35, BRANCO)
    add_rect(sl, lx, ty, 6.2, 0.07, VERDE)
    add_text(sl, tit, lx + 0.2, ty + 0.18, 5.8, 0.55,
             font_size=15, bold=True, color=VERDE)
    add_text(sl, desc, lx + 0.2, ty + 0.75, 5.8, 1.4,
             font_size=13, color=CINZA_MED)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 – Tecnologia (Resumo)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, BRANCO)
slide_header(sl, "Tecnologia", "Arquitectura moderna, segura e escalável")

tech_items = [
    ("⚡  Backend",          "FastAPI (Python) + PostgreSQL — rápido, robusto e open-source"),
    ("🌐  Web App",          "React 19 + Vite — interface moderna e responsiva"),
    ("📱  App Mobile",       "React Native + Expo — iOS e Android sem código duplicado"),
    ("🗺️  Mapas",            "Leaflet com rotação por bússola — mapas precisos e leves"),
    ("🔴  Tempo Real",       "WebSocket nativo — actualizações instantâneas sem reload"),
    ("💳  Pagamentos",       "Stripe — padrão mundial, seguro e certificado PCI"),
    ("☁️  Alojamento",       "Heroku/Render — infraestrutura cloud com escalabilidade automática"),
    ("🔒  Segurança",        "JWT + bcrypt — autenticação e passwords encriptadas"),
]

for i, (tit, desc) in enumerate(tech_items):
    row = i // 2
    col = i % 2
    lx = 0.4 + col * 6.45
    ty = 1.65 + row * 1.4
    add_rect(sl, lx, ty, 6.2, 1.25, AZUL_CLA)
    add_text(sl, tit, lx + 0.2, ty + 0.1, 5.8, 0.42,
             font_size=14, bold=True, color=AZUL_MAR)
    add_text(sl, desc, lx + 0.2, ty + 0.55, 5.8, 0.6,
             font_size=13, color=CINZA_MED)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 – Próximos Passos / CTA
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, AZUL_MAR)
add_rect(sl, 0, 5.5, 13.33, 2.0, AMARELO)
add_rect(sl, 0, 5.4, 13.33, 0.12, AMARELO_ESC)

add_text(sl, "Prontos para avançar?", 1.0, 0.9, 11.0, 0.9,
         font_size=38, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
add_text(sl, "Próximos Passos", 1.0, 1.75, 11.0, 0.55,
         font_size=22, bold=False, color=BRANCO, align=PP_ALIGN.CENTER)

steps_cta = [
    ("01", "Demo ao Vivo", "Agendamos uma demonstração prática com a câmara e convidados."),
    ("02", "Projecto Piloto", "Selecção de uma praia para teste da temporada. Implementação gratuita."),
    ("03", "Avaliação & Expansão", "Relatório de resultados ao fim da época. Decisão de alargamento."),
]

for i, (num, tit, desc) in enumerate(steps_cta):
    lx = 0.8 + i * 4.0
    add_rect(sl, lx, 2.55, 3.6, 2.6, RGBColor(0x00, 0x4F, 0x85))
    add_text(sl, num, lx + 0.15, 2.65, 0.8, 0.6,
             font_size=28, bold=True, color=AMARELO)
    add_text(sl, tit, lx + 0.15, 3.3, 3.3, 0.5,
             font_size=16, bold=True, color=BRANCO)
    add_text(sl, desc, lx + 0.15, 3.85, 3.3, 1.1,
             font_size=13, color=RGBColor(0xCC, 0xE5, 0xFF))

add_text(sl, "📧  suporte@sunnysales.com", 3.5, 5.62, 6.3, 0.55,
         font_size=19, bold=True, color=AZUL_MAR, align=PP_ALIGN.CENTER)
add_text(sl, "Transforme a sua praia numa praia inteligente.",
         2.0, 6.3, 9.33, 0.55,
         font_size=17, bold=False, color=AZUL_MAR, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
#  Guardar
# ════════════════════════════════════════════════════════════════════════════
output_path = "/home/user/ss_TESTER/SunnySales_Camara_Municipal.pptx"
prs.save(output_path)
print(f"Apresentação guardada em: {output_path}")
